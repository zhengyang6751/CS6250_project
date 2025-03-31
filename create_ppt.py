from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Sarcasm Detection Project"
    subtitle.text = "Using BERT for Text Classification"

def create_bullet_slide(prs, title, bullet_points):
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    
    title_shape = shapes.title
    title_shape.text = title
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

def main():
    prs = Presentation()
    
    # Title Slide
    create_title_slide(prs)
    
    # Project Overview
    create_bullet_slide(prs, "Project Overview", [
        "Goal: Detect sarcasm in text using machine learning",
        "Model: BERT (Bidirectional Encoder Representations from Transformers)",
        "Task: Binary classification (sarcastic vs non-sarcastic)"
    ])
    
    # Technical Stack
    create_bullet_slide(prs, "Technical Stack", [
        "Framework: PyTorch",
        "Model: BERT-base-uncased",
        "Libraries:",
        "• pandas: Data manipulation",
        "• transformers: BERT implementation",
        "• scikit-learn: Evaluation metrics",
        "• torch: Deep learning framework"
    ])
    
    # Data Processing
    create_bullet_slide(prs, "Data Processing", [
        "Input: Text comments",
        "Preprocessing:",
        "• Remove special characters",
        "• Convert to lowercase",
        "• Handle missing values",
        "Labels: Binary (contains_slash_s)"
    ])
    
    # Model Architecture
    create_bullet_slide(prs, "Model Architecture", [
        "Base Model: BERT-base-uncased",
        "Classification Head:",
        "• Input: 768-dimensional BERT embeddings",
        "• Output: 2 classes (sarcastic/non-sarcastic)",
        "Training Parameters:",
        "• Batch size: 16",
        "• Learning rate: 1e-5",
        "• Epochs: 3"
    ])
    
    # Training Process
    create_bullet_slide(prs, "Training Process", [
        "1. Data loading and preprocessing",
        "2. Tokenization using BERT tokenizer",
        "3. Model training with Adam optimizer",
        "4. Loss calculation and backpropagation",
        "5. Model evaluation"
    ])
    
    # Evaluation Metrics
    create_bullet_slide(prs, "Evaluation Metrics", [
        "Accuracy",
        "Precision",
        "Recall",
        "F1-score"
    ])
    
    # Future Improvements
    create_bullet_slide(prs, "Future Improvements", [
        "Try different BERT variants",
        "Experiment with different hyperparameters",
        "Add more preprocessing steps",
        "Implement cross-validation",
        "Add data augmentation"
    ])
    
    # Thank You Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Questions?"
    
    # Save the presentation
    prs.save('sarcasm_project.pptx')

if __name__ == '__main__':
    main() 